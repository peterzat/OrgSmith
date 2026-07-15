"""Unit tier: briefing decks (.pptx) end to end.

Planner quota and failure, airlock authoring, rendered deck opens in
python-pptx with facts, mentions, and the OPC marker, and the validator
catches a corrupted deck.
"""

import io
import shutil
import zipfile

import pytest

from orgsmith.artifacts import load_engagements, load_manifest
from orgsmith.assemble import run_assemble
from orgsmith.paths import OrgPaths
from orgsmith.render import run_render
from orgsmith.render.provenance import opc_has_marker
from orgsmith.validate import run_validate

from conftest import build_mix_stages, run_authoring, run_enrichment, write_mix_recipe

pytestmark = pytest.mark.unit

DECK_MIX = {"docx": 8, "pdf": 3, "xlsx": 2, "pptx": 2}


@pytest.fixture(scope="module")
def deck_org(tmp_path_factory):
    paths = build_mix_stages(tmp_path_factory.mktemp("deck-org"), DECK_MIX)
    run_enrichment(paths)
    run_authoring(paths)
    assert run_render(paths) == 0
    assert run_assemble(paths) == 0
    return paths


@pytest.fixture()
def deck_org_copy(deck_org, tmp_path):
    shutil.copytree(deck_org.root / "recipes", tmp_path / "recipes")
    shutil.copytree(deck_org.root / "companies", tmp_path / "companies")
    return OrgPaths(root=tmp_path, slug=deck_org.slug)


def _decks(paths):
    return [e for e in load_manifest(paths) if e.format == "pptx"]


def _pptx_text(path) -> str:
    from pptx import Presentation

    chunks = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.extend(p.text for p in shape.text_frame.paragraphs)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.extend(c.text for c in row.cells)
    return " ".join(chunks)


def test_decks_planned_per_engagement(deck_org):
    decks = _decks(deck_org)
    assert len(decks) == DECK_MIX["pptx"]
    engagements = {e.engagement for e in decks}
    assert len(engagements) == len(decks)  # one deck per engagement
    for deck in decks:
        assert deck.genre == "briefing_deck"
        assert deck.path.endswith(".pptx")
        assert deck.authoring == "batchable"
        assert deck.facts_refs, "deck plants no facts"
        assert deck.mentions, "deck plans no mentions"


def test_overdemanding_deck_mix_fails_actionably(tmp_path):
    from orgsmith.charter import run_charter
    from orgsmith.docplan import run_docplan
    from orgsmith.fabric import run_fabric
    from orgsmith.foundation import run_scaffold

    # 4 decks demanded, 3 engagements available.
    paths = write_mix_recipe(
        tmp_path, {"docx": 8, "pdf": 3, "xlsx": 2, "pptx": 4}
    )
    assert run_charter(paths) == 0
    assert run_scaffold(paths) == 0
    assert run_fabric(paths) == 0
    with pytest.raises(SystemExit, match="format_mix.pptx"):
        run_docplan(paths)


def test_rendered_deck_opens_with_facts_mentions_and_marker(deck_org):
    facts = load_engagements(deck_org).fact_index()
    for deck in _decks(deck_org):
        path = deck_org.share_dir / deck.path
        assert path.exists()
        assert opc_has_marker(path)
        text = _pptx_text(path)
        for ref in deck.facts_refs:
            assert facts[ref].rendered in text, f"{deck.path}: {ref}"
        for mention in deck.mentions:
            assert mention.surface in text, f"{deck.path}: {mention.surface}"


def test_deck_org_validates_clean(deck_org):
    assert run_validate(deck_org) == 0


def test_gutted_deck_fails_fact_mention_and_marker_rules(deck_org_copy, capsys):
    from pptx import Presentation

    deck = _decks(deck_org_copy)[0]
    blank = io.BytesIO()
    Presentation().save(blank)  # no slides, no marker, no text
    (deck_org_copy.share_dir / deck.path).write_bytes(blank.getvalue())
    assert run_validate(deck_org_copy) == 1
    out = capsys.readouterr().out
    for rule in ("FACT-01", "MENT-01", "PROV-01"):
        assert rule in out, rule
    assert deck.path.rsplit("/", 1)[-1] in out


def test_marker_stripped_deck_fails_prov(deck_org_copy, capsys):
    deck = _decks(deck_org_copy)[0]
    target = deck_org_copy.share_dir / deck.path
    src = zipfile.ZipFile(io.BytesIO(target.read_bytes()))
    out_buf = io.BytesIO()
    with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as out:
        for name in src.namelist():
            if name != "docProps/custom.xml":
                out.writestr(name, src.read(name))
    target.write_bytes(out_buf.getvalue())
    assert run_validate(deck_org_copy, only=["PROV-01"]) == 1
    assert "synthetic-provenance marker missing" in capsys.readouterr().out
