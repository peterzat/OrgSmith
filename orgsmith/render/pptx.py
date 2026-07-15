""".pptx renderer: python-pptx slides from DocIR, house style, core
properties, and the OPC synthetic marker.

Slide mapping: every heading block starts a new slide (its text is the
slide title); the paragraph/list/table blocks that follow fill that
slide's body. Blocks arriving before any heading open an untitled lead
slide named after the document.
"""

from __future__ import annotations

import io
from datetime import datetime, timezone

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ..fabric.engagements import render_date
from ..schemas import DocIR, ManifestEntry
from .provenance import add_opc_marker
from .styles import StylePack

_TITLE_LAYOUT = 0  # default-template "Title Slide"
_BODY_LAYOUT = 1  # default-template "Title and Content"


def _body_frame(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            return shape.text_frame
    return None


def _add_line(frame, text: str, bullet: bool) -> None:
    para = frame.paragraphs[0] if not frame.paragraphs[0].runs else frame.add_paragraph()
    para.text = text
    para.level = 1 if bullet else 0


def render_pptx(
    docir: DocIR,
    entry: ManifestEntry,
    style: StylePack,
    author_name: str,
) -> bytes:
    prs = Presentation()
    when_text = render_date(entry.date)
    slide = None
    frame = None
    first = True

    def new_slide(title_text: str):
        nonlocal slide, frame, first
        layout = prs.slide_layouts[_TITLE_LAYOUT if first else _BODY_LAYOUT]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = RGBColor.from_string(style.accent_hex)
        frame = _body_frame(slide)
        if first and frame is not None:
            # The title layout's subtitle doubles as the deck's dateline.
            frame.text = f"{style.letterhead_lines[0]} — {when_text}"
        first = False

    for block in docir.blocks:
        if block.kind == "heading":
            new_slide(block.text)
            continue
        if slide is None:
            new_slide(entry.title)
        if block.kind == "paragraph":
            if frame is not None:
                _add_line(frame, block.text, bullet=False)
        elif block.kind == "list":
            if frame is not None:
                for item in block.items:
                    _add_line(frame, item, bullet=True)
        elif block.kind == "table":
            rows = ([block.header] if block.header else []) + block.rows
            cols = max((len(r) for r in rows), default=1)
            shape = slide.shapes.add_table(
                len(rows), cols, Inches(0.6), Inches(2.2), Inches(8.8),
                Inches(0.4 * len(rows)),
            )
            for r, row in enumerate(rows):
                for c, text in enumerate(row):
                    cell = shape.table.cell(r, c)
                    cell.text = text
                    for run in cell.text_frame.paragraphs[0].runs:
                        run.font.size = Pt(12)
        elif block.kind == "sigblock":
            raise SystemExit(
                f"render: sigblock in {entry.doc_id} ({entry.format}); decks "
                "carry no signatures"
            )

    props = prs.core_properties
    props.title = entry.title
    props.author = author_name
    props.last_modified_by = author_name
    stamp = datetime(
        entry.date.year, entry.date.month, entry.date.day, 9, 0, 0,
        tzinfo=timezone.utc,
    )
    props.created = stamp
    props.modified = stamp

    buf = io.BytesIO()
    prs.save(buf)
    return add_opc_marker(buf.getvalue())
