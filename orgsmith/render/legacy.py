"""Legacy conversion: verified modern intermediates become .doc/.xls/.ppt.

Generation-time only, and the single place in the pipeline that shells
out: LibreOffice (`soffice --headless --convert-to`) turns a fully
rendered, marker-verified OPC intermediate into the pre-2007 binary the
manifest planned. Validation never comes back here; it reads the
binaries with olefile/xlrd, pure Python.

Marker strategy (see the M5 spike): the intermediate carries the normal
OPC custom-property marker, plus a `<marker>:true` token planted in the
core dc:description. LibreOffice drops custom properties in conversion
but maps dc:description onto OLE SummaryInformation comments, which is
what `legacy_has_marker` reads back. Render verifies the converted file
with that same checker and fails loudly when the marker did not survive.
"""

from __future__ import annotations

import io
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path

from ..schemas import Fact, ManifestEntry
from .provenance import MARKER_NAME, legacy_has_marker, opc_has_marker

LEGACY_STREAMS = {
    "doc": "WordDocument",
    "xls": "Workbook",
    "ppt": "PowerPoint Document",
}

_DESCRIPTION = re.compile(
    r"<dc:description\s*/>|<dc:description>.*?</dc:description>", re.DOTALL
)

SOFFICE_HINT = (
    "render: legacy documents are pending but `soffice` is not on PATH; "
    "install LibreOffice (libreoffice-writer/-calc/-impress) on the "
    "generation machine and confirm with `python -m orgsmith doctor`"
)


def require_soffice() -> str:
    soffice = shutil.which("soffice")
    if soffice is None:
        raise SystemExit(SOFFICE_HINT)
    return soffice


def _with_comments_token(opc_bytes: bytes) -> bytes:
    """Plant `<marker>:true` in docProps/core.xml dc:description, the one
    channel LibreOffice carries into the OLE binary."""
    src = zipfile.ZipFile(io.BytesIO(opc_bytes))
    buf = io.BytesIO()
    token = f"{MARKER_NAME}:true"
    planted = False
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as out:
        for name in src.namelist():
            payload = src.read(name)
            if name == "docProps/core.xml":
                text = _DESCRIPTION.sub("", payload.decode("utf-8"))
                close = "</cp:coreProperties>"
                if close not in text:
                    raise SystemExit(
                        "render: intermediate core.xml has no coreProperties "
                        "close tag; cannot plant the legacy marker token"
                    )
                text = text.replace(
                    close, f"<dc:description>{token}</dc:description>{close}"
                )
                payload = text.encode("utf-8")
                planted = True
            out.writestr(name, payload)
    if not planted:
        raise SystemExit(
            "render: intermediate has no docProps/core.xml; cannot plant "
            "the legacy marker token"
        )
    return buf.getvalue()


def _intermediate_text(data: bytes, fmt: str) -> str:
    if fmt == "doc":
        import docx

        d = docx.Document(io.BytesIO(data))
        chunks = [p.text for p in d.paragraphs]
        for t in d.tables:
            for row in t.rows:
                chunks.extend(c.text for c in row.cells)
        return "\n".join(chunks)
    if fmt == "ppt":
        from pptx import Presentation

        chunks = []
        for slide in Presentation(io.BytesIO(data)).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.extend(p.text for p in shape.text_frame.paragraphs)
                if shape.has_table:
                    for row in shape.table.rows:
                        chunks.extend(c.text for c in row.cells)
        return "\n".join(chunks)
    return ""  # xls: workbook values are FIN-02's business, not prose


def render_legacy(
    entry: ManifestEntry,
    intermediate: bytes,
    target: Path,
    facts: dict[str, Fact],
) -> None:
    """Verify the modern intermediate, convert it, verify the binary."""
    soffice = require_soffice()

    with tempfile.TemporaryDirectory(prefix="orgsmith-legacy-") as tmp:
        tmpdir = Path(tmp)
        base = "intermediate." + {"doc": "docx", "xls": "xlsx", "ppt": "pptx"}[
            entry.format
        ]
        src = tmpdir / base
        src.write_bytes(_with_comments_token(intermediate))

        if not opc_has_marker(src):
            raise SystemExit(
                f"render: modern intermediate for {entry.doc_id} lost its "
                "OPC marker before conversion"
            )
        text = re.sub(r"\s+", " ", _intermediate_text(intermediate, entry.format))
        for ref in entry.facts_refs:
            if ref in facts and facts[ref].rendered not in text:
                raise SystemExit(
                    f"render: fact {ref} missing from the modern intermediate "
                    f"of {entry.doc_id}; refusing to convert an unverified doc"
                )

        profile = tmpdir / "lo-profile"
        result = subprocess.run(
            [
                soffice,
                "--headless",
                f"-env:UserInstallation={profile.as_uri()}",
                "--convert-to",
                entry.format,
                "--outdir",
                str(tmpdir),
                str(src),
            ],
            capture_output=True,
            timeout=300,
        )
        converted = tmpdir / f"intermediate.{entry.format}"
        if result.returncode != 0 or not converted.exists():
            raise SystemExit(
                f"render: soffice failed converting {entry.doc_id} to "
                f"{entry.format}: {result.stderr.decode('utf-8', 'replace')}"
            )
        target.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(converted), str(target))

    if not legacy_has_marker(target):
        raise SystemExit(
            f"render: converted {target.name} carries no synthetic marker "
            "in its OLE comments; refusing to ship an unmarked binary"
        )
