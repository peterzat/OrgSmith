"""Extractable text of rendered documents: the one shared reader.

Three consumers have to agree exactly on "what text does this file actually
expose?": the validator's text obligations (FACT-01, MENT-01, LOC-01/02),
the eval emitter's rendered-truth scans, and the retrieval baselines. This
module is that single answer, so they cannot drift apart. It was lifted out
of `validate/rules.py` unchanged; the validator `Context` now delegates to
it.

Pure-python readers only. Nothing here imports weasyprint at module scope,
so a validate-only or score-only install never needs the Pango stack, and
nothing here touches the network.
"""

from __future__ import annotations

import re


def mask_surfaces(text: str, surfaces) -> str:
    """Text with the given surfaces removed, so a scan for a short token
    cannot match inside a longer planned name (an alias `Jim` standing
    inside a planned `Jim Halpert` belongs to Halpert, not to whoever
    registered `Jim`).

    Public because it is a shared contract, not one module's helper: the
    eval emitter's alias diagnostics and the MENT-03 validator rule must
    mask identically or they disagree about the same document. Longest
    surface first, and the same non-word lookarounds `surface_in_text`
    uses, so masking and scanning agree on what a standalone token is."""
    for surface in sorted(surfaces, key=len, reverse=True):
        if surface:
            text = re.sub(rf"(?<!\w){re.escape(surface)}(?!\w)", " ", text)
    return text


def image_only(entry) -> bool:
    """A scan with no OCR layer exposes no extractable text by design."""
    return bool(entry.render_params.get("scan")) and not entry.render_params.get(
        "ocr_layer"
    )


def pptx_text(path) -> str:
    """Every text run a pptx exposes: shape frames plus table cells."""
    from pptx import Presentation

    chunks: list[str] = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.extend(p.text for p in shape.text_frame.paragraphs)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.extend(c.text for c in row.cells)
    return "\n".join(chunks)


def eml_message(path):
    from email import policy
    from email.parser import BytesParser

    with open(path, "rb") as fh:
        return BytesParser(policy=policy.default).parse(fh)


class DocText:
    """Cached extractable text for one org's rendered documents.

    `engagements` and `foundation` are read only by the legacy
    (`.doc`/`.ppt`) path, whose text obligations run against the authoring
    DocIR rather than a binary-format parser; pass them whenever the org
    may hold legacy binaries. One instance per org keeps the extraction
    cost paid once no matter how many consumers ask.
    """

    def __init__(self, paths, engagements=None, foundation=None):
        self.paths = paths
        self.engagements = engagements
        self.foundation = foundation
        self._cache: dict = {}

    def text(self, entry) -> str:
        """Extractable text of a rendered doc, whitespace-normalized."""
        if entry.doc_id in self._cache:
            return self._cache[entry.doc_id]
        path = self.paths.share_dir / entry.path
        if entry.format == "docx":
            import docx

            d = docx.Document(str(path))
            chunks = [p.text for p in d.paragraphs]
            for t in d.tables:
                for row in t.rows:
                    chunks.extend(c.text for c in row.cells)
            text = "\n".join(chunks)
        elif entry.format == "pdf":
            if image_only(entry):
                # An image-only scan exposes no text by design; its text
                # obligations run against the archived truth. A missing or
                # unreadable archive yields empty text here, so FACT/MENT
                # fail loudly alongside SCAN-02, never a silent pass.
                text = "\n".join(self.scan_archive(entry) or [])
            else:
                from pypdf import PdfReader

                text = "\n".join(
                    page.extract_text() or ""
                    for page in PdfReader(str(path)).pages
                )
        elif entry.format == "pptx":
            text = pptx_text(path)
        elif entry.format == "eml":
            from email.utils import getaddresses

            msg = eml_message(path)
            body = msg.get_body(preferencelist=("plain",))
            body_text = body.get_content() if body is not None else ""
            # Include the To/Cc header display names: a recipient's full name
            # legitimately lives in the transport headers, not the body, so
            # MENT-01 must read it there. The eml body no longer carries a
            # To:/Cc: banner (strip_leading_header_block), so without this a
            # header-only recipient name would read as a missing mention.
            names = [
                nm
                for nm, _addr in getaddresses(
                    msg.get_all("To", []) + msg.get_all("Cc", [])
                )
                if nm
            ]
            text = (" ".join(names) + "\n" + body_text) if names else body_text
        elif entry.format in ("xlsx", "xls"):
            # Workbooks are checked cell-by-cell (FIN-02), not as prose;
            # FACT-01 skips them explicitly.
            text = ""
        elif entry.format in ("doc", "ppt"):
            text = self._legacy_text(entry)
        else:
            raise SystemExit(
                f"validate: no text extractor for format {entry.format!r} "
                f"({entry.path})"
            )
        text = re.sub(r"\s+", " ", text)
        self._cache[entry.doc_id] = text
        return text

    def pages(self, entry) -> list[str]:
        """Per-page extractable text (pdf only), whitespace-normalized.
        Page addressing is what makes signature-page scoping checkable.
        Image-only scans page-address the archived truth instead."""
        key = ("pages", entry.doc_id)
        if key in self._cache:
            return self._cache[key]
        if image_only(entry):
            raw = self.scan_archive(entry) or []
        else:
            from pypdf import PdfReader

            path = self.paths.share_dir / entry.path
            raw = [
                page.extract_text() or "" for page in PdfReader(str(path)).pages
            ]
        pages = [re.sub(r"\s+", " ", p) for p in raw]
        self._cache[key] = pages
        return pages

    def layout_text(self, entry) -> str:
        """PDF text via pypdf's 'layout' extraction, whitespace-normalized.

        The default 'plain' mode inserts a spurious intra-word space for some
        glyph sequences (a hyphen before a capital: "Kirby-Taylor" comes out
        "Kirby-T aylor"), a false negative for a surface-in-text check even
        though the name renders correctly. Layout mode preserves the run.

        Used ONLY as a FACT-01/MENT-01 fallback when the surface is missing
        from the plain text, so it can rescue a rendered-but-mis-extracted
        surface without changing any check that already passes. Returns "" for
        non-pdf or image-only entries, so the fallback is a no-op there."""
        key = ("layout", entry.doc_id)
        if key in self._cache:
            return self._cache[key]
        if entry.format != "pdf" or image_only(entry):
            self._cache[key] = ""
            return ""
        from pypdf import PdfReader

        path = self.paths.share_dir / entry.path
        text = "\n".join(
            page.extract_text(extraction_mode="layout") or ""
            for page in PdfReader(str(path)).pages
        )
        text = re.sub(r"\s+", " ", text)
        self._cache[key] = text
        return text

    def _legacy_text(self, entry) -> str:
        """Text obligations for a converted binary run against its
        authoring source: the fact-resolved DocIR (plus signer names, which
        the modern renderers print as signature lines). Reading prose back
        out of .doc/.ppt would need a binary-format parser; conversion
        fidelity is a documented residual risk, but the DocIR is exactly
        what the verified modern intermediate rendered. Missing or
        unresolvable DocIR yields empty text, so FACT/MENT fail loudly."""
        from .authoring.ingest import docir_path
        from .render import people_index
        from .render.resolve import FactResolutionError, resolve_docir
        from .schemas import DocIR

        source = docir_path(self.paths, entry.doc_id)
        if not source.exists():
            return ""
        try:
            resolved = resolve_docir(
                DocIR.model_validate_json(source.read_text("utf-8")),
                self.engagements.fact_index(),
            )
        except FactResolutionError:
            return ""
        people = people_index(self.foundation)
        chunks: list[str] = []
        for b in resolved.blocks:
            chunks.append(b.text)
            chunks.extend(b.items)
            chunks.extend(b.header)
            for row in b.rows:
                chunks.extend(row)
            for signer in b.signers:
                if signer in people:
                    chunks.append(people[signer]["name"])
        return "\n".join(chunks)

    def scan_archive(self, entry) -> list | None:
        """Archived true per-page text for a scanned doc, or None when the
        archive is missing or does not parse (SCAN-02's findings)."""
        key = ("archive", entry.doc_id)
        if key not in self._cache:
            from .render.scan import scan_pages_path
            from .schemas import ScanPages

            path = scan_pages_path(self.paths, entry.doc_id)
            pages = None
            if path.exists():
                try:
                    pages = ScanPages.model_validate_json(
                        path.read_text("utf-8")
                    ).pages
                except Exception:  # noqa: BLE001 - unparseable = absent
                    pages = None
            self._cache[key] = pages
        return self._cache[key]
