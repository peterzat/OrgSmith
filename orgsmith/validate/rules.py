"""The v0 rule catalog: ORG, DATE, FIN, FACT, FILE, MAN, PROV families."""

from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass, field
from typing import Callable, Iterator

from ..artifacts import (
    load_acl,
    load_charter,
    load_distribution_lists,
    load_engagements,
    load_finance,
    load_foundation,
    load_graph,
    load_manifest,
    load_mention_map,
)
from ..paths import OrgPaths
from ..schemas import surface_in_text

Finding = tuple[str, str]  # (message, target)


@dataclass
class Context:
    paths: OrgPaths
    charter: object
    foundation: object
    finance: object
    engagements: object
    graph: object
    mention_map: object  # None for orgs generated before mention ground truth
    acl: object  # None for orgs generated before the ACL overlay
    manifest: list
    dls: object = None  # None when the recipe declares no distribution lists
    _text_cache: dict = field(default_factory=dict)

    @classmethod
    def load(cls, paths: OrgPaths) -> "Context":
        return cls(
            paths=paths,
            charter=load_charter(paths),
            foundation=load_foundation(paths),
            finance=load_finance(paths),
            engagements=load_engagements(paths),
            graph=load_graph(paths),
            mention_map=load_mention_map(paths),
            acl=load_acl(paths),
            manifest=load_manifest(paths),
            dls=load_distribution_lists(paths),
        )

    def entry(self, doc_id: str):
        for e in self.manifest:
            if e.doc_id == doc_id:
                return e
        return None

    def doc_text(self, entry) -> str:
        """Extractable text of a rendered doc, whitespace-normalized."""
        if entry.doc_id in self._text_cache:
            return self._text_cache[entry.doc_id]
        path = self.paths.share_dir / entry.path
        if entry.format == "docx":
            import docx

            d = docx.Document(str(path))
            chunks = [p.text for p in d.paragraphs]
            for t in d.tables:
                for row in t.rows:
                    chunks.extend(c.text for c in row.cells)
            text = "\n".join(chunks)
        elif entry.format == "pdf":
            if _image_only(entry):
                # An image-only scan exposes no text by design; its text
                # obligations run against the archived truth. A missing or
                # unreadable archive yields empty text here, so FACT/MENT
                # fail loudly alongside SCAN-02, never a silent pass.
                text = "\n".join(self.scan_archive(entry) or [])
            else:
                from pypdf import PdfReader

                text = "\n".join(
                    page.extract_text() or ""
                    for page in PdfReader(str(path)).pages
                )
        elif entry.format == "pptx":
            text = _pptx_text(path)
        elif entry.format == "eml":
            body = _eml_message(path).get_body(preferencelist=("plain",))
            text = body.get_content() if body is not None else ""
        elif entry.format in ("xlsx", "xls"):
            # Workbooks are checked cell-by-cell (FIN-02), not as prose;
            # FACT-01 skips them explicitly.
            text = ""
        elif entry.format in ("doc", "ppt"):
            text = self._legacy_text(entry)
        else:
            raise SystemExit(
                f"validate: no text extractor for format {entry.format!r} "
                f"({entry.path})"
            )
        text = re.sub(r"\s+", " ", text)
        self._text_cache[entry.doc_id] = text
        return text

    def doc_pages(self, entry) -> list[str]:
        """Per-page extractable text (pdf only), whitespace-normalized.
        Page addressing is what makes signature-page scoping checkable.
        Image-only scans page-address the archived truth instead."""
        key = ("pages", entry.doc_id)
        if key in self._text_cache:
            return self._text_cache[key]
        if _image_only(entry):
            raw = self.scan_archive(entry) or []
        else:
            from pypdf import PdfReader

            path = self.paths.share_dir / entry.path
            raw = [
                page.extract_text() or "" for page in PdfReader(str(path)).pages
            ]
        pages = [re.sub(r"\s+", " ", p) for p in raw]
        self._text_cache[key] = pages
        return pages

    def pdf_layout_text(self, entry) -> str:
        """PDF text via pypdf's 'layout' extraction, whitespace-normalized.

        The default 'plain' mode inserts a spurious intra-word space for some
        glyph sequences (a hyphen before a capital: "Kirby-Taylor" comes out
        "Kirby-T aylor"), a false negative for a surface-in-text check even
        though the name renders correctly. Layout mode preserves the run.

        Used ONLY as a FACT-01/MENT-01 fallback when the surface is missing
        from the plain text, so it can rescue a rendered-but-mis-extracted
        surface without changing any check that already passes. Returns "" for
        non-pdf or image-only entries, so the fallback is a no-op there."""
        key = ("layout", entry.doc_id)
        if key in self._text_cache:
            return self._text_cache[key]
        if entry.format != "pdf" or _image_only(entry):
            self._text_cache[key] = ""
            return ""
        from pypdf import PdfReader

        path = self.paths.share_dir / entry.path
        text = "\n".join(
            page.extract_text(extraction_mode="layout") or ""
            for page in PdfReader(str(path)).pages
        )
        text = re.sub(r"\s+", " ", text)
        self._text_cache[key] = text
        return text

    def _legacy_text(self, entry) -> str:
        """Text obligations for a converted binary run against its
        authoring source: the fact-resolved DocIR (plus signer names, which
        the modern renderers print as signature lines). Reading prose back
        out of .doc/.ppt would need a binary-format parser; conversion
        fidelity is a documented residual risk, but the DocIR is exactly
        what the verified modern intermediate rendered. Missing or
        unresolvable DocIR yields empty text, so FACT/MENT fail loudly."""
        from ..authoring.ingest import docir_path
        from ..render import people_index
        from ..render.resolve import FactResolutionError, resolve_docir
        from ..schemas import DocIR

        source = docir_path(self.paths, entry.doc_id)
        if not source.exists():
            return ""
        try:
            resolved = resolve_docir(
                DocIR.model_validate_json(source.read_text("utf-8")),
                self.engagements.fact_index(),
            )
        except FactResolutionError:
            return ""
        people = people_index(self.foundation)
        chunks: list[str] = []
        for b in resolved.blocks:
            chunks.append(b.text)
            chunks.extend(b.items)
            chunks.extend(b.header)
            for row in b.rows:
                chunks.extend(row)
            for signer in b.signers:
                if signer in people:
                    chunks.append(people[signer]["name"])
        return "\n".join(chunks)

    def scan_archive(self, entry) -> list | None:
        """Archived true per-page text for a scanned doc, or None when the
        archive is missing or does not parse (SCAN-02's findings)."""
        key = ("archive", entry.doc_id)
        if key not in self._text_cache:
            from ..render.scan import scan_pages_path
            from ..schemas import ScanPages

            path = scan_pages_path(self.paths, entry.doc_id)
            pages = None
            if path.exists():
                try:
                    pages = ScanPages.model_validate_json(
                        path.read_text("utf-8")
                    ).pages
                except Exception:  # noqa: BLE001 - unparseable = absent
                    pages = None
            self._text_cache[key] = pages
        return self._text_cache[key]


def _pptx_text(path) -> str:
    """Every text run a pptx exposes: shape frames plus table cells."""
    from pptx import Presentation

    chunks: list[str] = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.extend(p.text for p in shape.text_frame.paragraphs)
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.extend(c.text for c in row.cells)
    return "\n".join(chunks)


def _eml_message(path):
    from email import policy
    from email.parser import BytesParser

    with open(path, "rb") as fh:
        return BytesParser(policy=policy.default).parse(fh)


def _image_only(entry) -> bool:
    return bool(entry.render_params.get("scan")) and not entry.render_params.get(
        "ocr_layer"
    )


def _always(_ctx: Context) -> str | None:
    return None


def _needs_eml(ctx: Context) -> str | None:
    mail = ctx.charter.doc_culture.mail
    mundane = mail.mundane_emails if mail is not None else 0
    if ctx.charter.doc_culture.format_mix.eml == 0 and mundane == 0:
        return "format_mix.eml is 0 and no mundane mail for this recipe"
    return None


def _needs_mail(ctx: Context) -> str | None:
    if ctx.charter.doc_culture.mail is None:
        return "doc_culture.mail is not declared for this recipe"
    return None


def _needs_dls(ctx: Context) -> str | None:
    mail = ctx.charter.doc_culture.mail
    if mail is None or mail.distribution_lists == 0:
        return "no distribution lists declared for this recipe"
    return None


def _needs_style(ctx: Context) -> str | None:
    if not ctx.charter.doc_culture.style_specs:
        return "recipe declares no per-person style specs"
    return None


def sty_01(ctx: Context):
    """M15: the style-spec ledger recomputes exactly from charter + roster
    via the one shared derive_style_specs twin. Grandfathers by charter:
    skips when style_specs is off; a knob on with the ledger missing or
    hand-edited is a failure, not a skip. The spec is never model-authored,
    so any drift from the recompute is tamper evidence."""
    from ..artifacts import load_style_specs
    from ..foundation.style import derive_style_specs

    got = load_style_specs(ctx.paths)
    if got is None:
        yield (
            "style_specs is on but ledger/style_specs.json is missing",
            "ledger/style_specs.json",
        )
        return
    want = derive_style_specs(ctx.charter, ctx.foundation)
    if got != want:
        yield (
            "style-spec ledger does not recompute from charter + roster",
            "ledger/style_specs.json",
        )


def _needs_scan(ctx: Context) -> str | None:
    if ctx.charter.doc_culture.scanned_ratio == 0:
        return "scanned_ratio is 0 for this recipe"
    return None


def _needs_legacy(ctx: Context) -> str | None:
    if ctx.charter.doc_culture.legacy_ratio == 0:
        return "legacy_ratio is 0 for this recipe"
    return None


def _needs_mentions(ctx: Context) -> str | None:
    if ctx.mention_map is None:
        return "org predates mention ground truth (no mention_map.json)"
    return None


def _needs_acl(ctx: Context) -> str | None:
    # Only an open posture may legitimately lack the ledger (every pre-ACL
    # org is posture open). A non-open org with no acl.json is corruption:
    # the rules run and yield the missing-ledger finding instead of skipping.
    if ctx.acl is None and ctx.charter.acl_posture == "open":
        return "org predates the ACL overlay (no ledger/acl.json)"
    return None


def _needs_mention_knob(ctx: Context) -> str | None:
    absent = _needs_mentions(ctx)
    if absent:
        return absent
    if ctx.charter.graph_targets.min_mentions_per_person < 1:
        return "min_mentions_per_person knob is 0 for this recipe"
    return None


def _needs_calendar(ctx: Context) -> str | None:
    if ctx.charter.doc_culture.business_calendar is None:
        return "business_calendar is not declared for this recipe"
    return None


def _needs_noise(ctx: Context) -> str | None:
    if ctx.charter.doc_culture.noise is None:
        return "noise model is not declared for this recipe"
    return None


@dataclass(frozen=True)
class Rule:
    id: str
    severity: str
    description: str
    check: Callable[[Context], Iterator[Finding]]
    # Returns a skip reason when the rule cannot run for this org (e.g. it
    # predates an artifact); None means run. Skips are surfaced visibly.
    available: Callable[[Context], str | None] = _always


def _employed_at(person, when) -> bool:
    emp = person.employment
    return emp.start <= when and (emp.end is None or emp.end >= when)


# --- ORG ------------------------------------------------------------------


def org_01(ctx: Context):
    roots = [p for p in ctx.foundation.people if p.reports_to is None]
    if len(roots) != 1:
        yield (
            f"expected exactly one CEO-equivalent, found {len(roots)}: "
            f"{[p.id for p in roots]}",
            "foundation.json",
        )


def org_02(ctx: Context):
    ids = {p.id for p in ctx.foundation.people}
    for p in ctx.foundation.people:
        seen = set()
        cur = p
        while cur.reports_to is not None:
            if cur.reports_to not in ids:
                yield (f"{cur.id} reports to unknown {cur.reports_to}",
                       "foundation.json")
                break
            if cur.id in seen:
                yield (f"reports_to cycle involving {cur.id}", "foundation.json")
                break
            seen.add(cur.id)
            cur = ctx.foundation.person(cur.reports_to)


# --- NAME -----------------------------------------------------------------


def name_01(ctx: Context):
    """Deliberately no grandfather: the screen reads only charter and
    foundation, artifacts every org has. The committed list is a screen,
    not a guarantee."""
    from ..namescreen import screen_charter, screen_foundation

    yield from screen_charter(ctx.charter)
    yield from screen_foundation(ctx.foundation)


# --- DATE -----------------------------------------------------------------


def date_01(ctx: Context):
    lo, hi = ctx.charter.doc_culture.date_range
    for e in ctx.manifest:
        if not (lo <= e.date <= hi):
            yield (f"doc date {e.date} outside charter range {lo}..{hi}", e.path)


def date_02(ctx: Context):
    for e in ctx.manifest:
        for author in e.authors:
            person = ctx.foundation.person(author)
            if not _employed_at(person, e.date):
                yield (
                    f"author {author} not employed on {e.date} "
                    f"(start {person.employment.start})",
                    e.path,
                )


def cal_01(ctx: Context):
    """M12: when a recipe declares a business-day calendar, every document of
    a genre that asserts a session happened (minutes, engagement mail) lands on
    a weekday that is not a declared holiday. The planner shifts these dates;
    this recomputes that no attendance document slipped onto a weekend or a
    declared holiday. Grandfathers by charter: skips visibly when the knob is
    off (available=_needs_calendar), and a knob that is on with a violation is
    a failure, not a skip."""
    from ..docplan.registry import REGISTRY

    holidays = set(ctx.charter.doc_culture.business_calendar.holidays)
    attendance = {r.genre for r in REGISTRY if r.asserts_attendance}
    for e in ctx.manifest:
        if e.genre not in attendance:
            continue
        if e.date.weekday() >= 5:
            yield (
                f"{e.genre} dated {e.date} ({e.date:%A}), a weekend; an "
                f"attendance genre must land on a business day",
                e.path,
            )
        elif e.date in holidays:
            yield (
                f"{e.genre} dated {e.date}, a declared holiday; an attendance "
                f"genre must land on a business day",
                e.path,
            )


def noise_01(ctx: Context):
    """M12: every derived noise document names a real authored source, so
    authored and derived words stay separable and a scorer can exclude noise.
    Grandfathers by charter: skips when noise is off (available=_needs_noise);
    a knob on with the noise stripped or a label broken is a failure, not a
    skip. File existence is MAN-01/FILE-01's job; this checks the labels --
    plus, since M15, the version-chain invariants: members cover 1..len-1
    with one agreed length, every member predates its final, and no two chain
    members (final included) are byte-identical, so a chain that hash dedupe
    could collapse is tamper evidence."""
    from ..render.noise import expected_empty_dirs

    noise = ctx.charter.doc_culture.noise
    for rel in expected_empty_dirs(ctx.charter, ctx.manifest):
        target = ctx.paths.share_dir / rel
        if not target.is_dir():
            yield (
                f"planned empty directory {rel!r} is missing from the share",
                rel,
            )
        elif next(target.iterdir(), None) is not None:
            yield (
                f"planned empty directory {rel!r} is not empty",
                rel,
            )
    by_id = {e.doc_id: e for e in ctx.manifest}
    derived = [e for e in ctx.manifest if e.authoring == "derived"]
    plans_files = any(
        (
            noise.duplicates,
            noise.drafts,
            noise.version_chains,
            noise.misfiled,
            noise.stale_templates,
        )
    )
    if not derived:
        if plans_files:
            yield (
                "noise model is declared but no derived documents were "
                "planted; the noise ground truth is missing",
                "docplan/manifest.jsonl",
            )
        return
    for e in derived:
        if e.noise_kind == "stale_template":
            # No source; the invariant is the content: a dead template's
            # every field is a bracketed dummy, so brackets must survive in
            # extractable text (a filled-in template is tamper).
            if (ctx.paths.share_dir / e.path).exists():
                text = ctx.doc_text(e)
                if "[" not in text or "]" not in text:
                    yield (
                        f"stale template {e.doc_id} carries no bracketed "
                        f"dummy field; a filled-in template is not a "
                        f"template",
                        e.path,
                    )
            continue
        src = by_id.get(e.noise_of)
        if src is None:
            yield (
                f"derived doc {e.doc_id} names source {e.noise_of!r}, which is "
                f"not in the manifest",
                e.path,
            )
        elif src.authoring == "derived":
            yield (
                f"derived doc {e.doc_id} derives from another derived doc "
                f"{e.noise_of}; noise must derive from an authored source",
                e.path,
            )
    for e in derived:
        if e.noise_kind != "misfile" or e.noise_of not in by_id:
            continue
        src = by_id[e.noise_of]
        e_dir = e.path.rsplit("/", 1)[0] if "/" in e.path else ""
        s_dir = src.path.rsplit("/", 1)[0] if "/" in src.path else ""
        if e_dir == s_dir:
            yield (
                f"misfiled doc {e.doc_id} sits in its source's own folder "
                f"{s_dir!r}; a misfile lives elsewhere",
                e.path,
            )
    chains: dict[str, list] = {}
    for e in derived:
        if e.noise_kind == "version" and e.noise_of in by_id:
            chains.setdefault(e.noise_of, []).append(e)
    for src_id, members in sorted(chains.items()):
        src = by_id[src_id]
        lens = {int(m.render_params["noise_len"]) for m in members}
        if len(lens) != 1:
            yield (
                f"version chain on {src_id} disagrees on noise_len "
                f"({sorted(lens)})",
                src.path,
            )
            continue
        length = lens.pop()
        positions = sorted(
            int(m.render_params["noise_pos"]) for m in members
        )
        if positions != list(range(1, length)):
            yield (
                f"version chain on {src_id} has members at {positions}, "
                f"wants 1..{length - 1}",
                src.path,
            )
        for m in members:
            if m.date >= src.date:
                yield (
                    f"version member {m.doc_id} is dated {m.date}, not "
                    f"earlier than its final {src_id} ({src.date})",
                    m.path,
                )
        digests: dict[str, list] = {}
        for m in [*members, src]:
            path = ctx.paths.share_dir / m.path
            if not path.exists():
                continue  # missing files are MAN-01/FILE-01 findings
            digest = hashlib.sha256(path.read_bytes()).hexdigest()
            digests.setdefault(digest, []).append(m.doc_id)
        for ids in digests.values():
            if len(ids) > 1:
                yield (
                    f"version chain on {src_id} carries byte-identical "
                    f"members {ids}; every version must diverge",
                    src.path,
                )
    by_path = {e.path: e for e in ctx.manifest}
    mismatched = sum(
        1
        for e in ctx.manifest
        if e.authoring != "derived"
        and e.render_params.get("attach_path")
        and by_path.get(str(e.render_params["attach_path"])) is not None
        and by_path[str(e.render_params["attach_path"])].noise_kind
        == "version"
    )
    if mismatched != noise.attachment_mismatch:
        yield (
            f"charter plans {noise.attachment_mismatch} attachment-version "
            f"mismatch(es) but {mismatched} transmittal(s) attach a "
            f"version-chain member",
            "docplan/manifest.jsonl",
        )


# --- FIN ------------------------------------------------------------------


def fin_01(ctx: Context):
    for check in ctx.finance.checks:
        if not check.ok:
            yield (f"ledger check failed: {check.name} ({check.detail})",
                   "ledger/finance.json")
    for fy in ctx.finance.years:
        if sum(fy.quarters) != fy.revenue:
            yield (
                f"FY{fy.year} quarters sum {sum(fy.quarters)} != revenue "
                f"{fy.revenue}",
                "ledger/finance.json",
            )


def fin_02(ctx: Context):
    for e in ctx.manifest:
        if e.format not in ("xlsx", "xls"):
            continue
        path = ctx.paths.share_dir / e.path
        if not path.exists():
            continue  # FILE-01 reports the absence
        year = int(e.render_params["year"])
        fy = next((y for y in ctx.finance.years if y.year == year), None)
        if fy is None:
            yield (f"workbook year {year} not in finance ledger", e.path)
            continue
        if e.format == "xlsx":
            from openpyxl import load_workbook

            cached = load_workbook(path, data_only=True)["Summary"]
            revenue = cached["F4"].value
            quarters = [cached[f"{c}4"].value for c in "BCDE"]
        else:
            # LibreOffice preserves cached formula results in .xls (spiked
            # at M5); xlrd reads them back without recalculation.
            import xlrd

            try:
                sheet = xlrd.open_workbook(str(path)).sheet_by_name("Summary")
            except Exception as err:  # noqa: BLE001 - unreadable = finding
                yield (f"xls workbook does not open via xlrd: {err}", e.path)
                continue
            revenue = sheet.cell_value(3, 5)
            quarters = [sheet.cell_value(3, c) for c in range(1, 5)]
        if revenue != fy.revenue:
            yield (
                f"cached revenue total {revenue} != ledger {fy.revenue}",
                e.path,
            )
        if quarters != fy.quarters:
            yield (f"cached quarters {quarters} != ledger {fy.quarters}", e.path)


# --- FACT (fact echo) ------------------------------------------------------


def fact_01(ctx: Context):
    facts = ctx.engagements.fact_index()
    for e in ctx.manifest:
        if not e.facts_refs or e.format == "xlsx":
            continue
        if not (ctx.paths.share_dir / e.path).exists():
            continue
        text = ctx.doc_text(e)
        for ref in e.facts_refs:
            if ref not in facts:
                yield (f"facts_ref {ref} not in engagement ledger", e.path)
            elif facts[ref].rendered not in text and (
                facts[ref].rendered not in ctx.pdf_layout_text(e)
            ):
                yield (
                    f"fact {ref} surface form {facts[ref].rendered!r} not "
                    f"found in extractable text",
                    e.path,
                )


# --- FILE -----------------------------------------------------------------


def file_01(ctx: Context):
    for e in ctx.manifest:
        path = ctx.paths.share_dir / e.path
        if not path.exists():
            yield ("manifest doc missing from share", e.path)
            continue
        try:
            if e.format == "docx":
                import docx

                docx.Document(str(path))
            elif e.format == "pdf":
                from pypdf import PdfReader

                if len(PdfReader(str(path)).pages) < 1:
                    yield ("pdf has no pages", e.path)
            elif e.format == "xlsx":
                from openpyxl import load_workbook

                load_workbook(path, data_only=True)
            elif e.format == "pptx":
                from pptx import Presentation

                if len(Presentation(str(path)).slides) < 1:
                    yield ("pptx has no slides", e.path)
            elif e.format == "eml":
                msg = _eml_message(path)
                for header in ("From", "To", "Date", "Subject", "Message-ID"):
                    if msg[header] is None:
                        yield (f"eml missing {header} header", e.path)
                # get_content() raises on a multipart transmittal; decode the
                # plain body and every attachment part instead (M14).
                body = msg.get_body(preferencelist=("plain",))
                if body is None:
                    yield ("eml has no plain-text body", e.path)
                else:
                    body.get_content()  # undecodable body -> reader failure
                for part in msg.iter_attachments():
                    part.get_content()  # undecodable attachment -> failure
            elif e.format in ("doc", "xls", "ppt"):
                import olefile

                if not olefile.isOleFile(str(path)):
                    yield ("legacy file is not an OLE container", e.path)
                else:
                    olefile.OleFileIO(str(path)).close()
            else:
                # Every supported format must have a reader branch above; a
                # format this rule does not know is a finding, not a pass.
                yield (f"no native reader known for format {e.format!r}", e.path)
        except Exception as err:  # noqa: BLE001 - any reader failure is the finding
            yield (f"file does not open in native reader: {err}", e.path)


# --- MAN ------------------------------------------------------------------

_SHARE_EXTRAS = {"TOC.md"}


def man_01(ctx: Context):
    if not ctx.paths.share_dir.exists():
        yield ("share directory missing", str(ctx.paths.share_dir))
        return
    on_disk = {
        str(p.relative_to(ctx.paths.share_dir))
        for p in ctx.paths.share_dir.rglob("*")
        if p.is_file()
    }
    planned = {e.path for e in ctx.manifest}
    # PERMISSIONS.md is a sanctioned extra only while the ACL ledger that
    # renders it exists; a stray copy on a ledgerless org is unmanifested.
    extras = _SHARE_EXTRAS | (
        {"PERMISSIONS.md"} if ctx.acl is not None else set()
    )
    for extra in sorted(on_disk - planned - extras):
        yield ("file in share but not in manifest", extra)
    for missing in sorted(planned - on_disk):
        yield ("manifest doc missing from share", missing)


# --- MENT (mention echo) ----------------------------------------------------


def ment_01(ctx: Context):
    for record in ctx.mention_map.mentions:
        entry = ctx.entry(record.doc_id)
        if entry is None:
            yield (f"mention references unknown doc {record.doc_id}",
                   "ledger/mention_map.json")
            continue
        if not (ctx.paths.share_dir / entry.path).exists():
            continue  # FILE-01/MAN-01 report the absence
        if not surface_in_text(
            record.surface, ctx.doc_text(entry)
        ) and not surface_in_text(
            record.surface, ctx.pdf_layout_text(entry)
        ):
            yield (
                f"planned mention surface {record.surface!r} "
                f"({record.entity}) not found in extractable text",
                entry.path,
            )


def ment_02(ctx: Context):
    entities = set(ctx.graph.entities)
    for record in ctx.mention_map.mentions:
        if record.entity not in entities:
            yield (
                f"mention entity {record.entity} not in the graph ledger",
                "ledger/mention_map.json",
            )


# --- GRAPH ------------------------------------------------------------------


def graph_01(ctx: Context):
    minimum = ctx.charter.graph_targets.min_mentions_per_person
    for person in ctx.foundation.people:
        docs = {
            r.doc_id for r in ctx.mention_map.mentions if r.entity == person.id
        }
        if len(docs) < minimum:
            yield (
                f"{person.id} planned in {len(docs)} docs, recipe requires "
                f">= {minimum}",
                "ledger/mention_map.json",
            )


def graph_02(ctx: Context):
    mentioned = {r.entity for r in ctx.mention_map.mentions}
    for person in ctx.foundation.people:
        if person.id not in mentioned:
            yield (f"roster member {person.id} has zero planned mentions",
                   "ledger/mention_map.json")


def graph_03(ctx: Context):
    valid = set(ctx.graph.entities) | {
        e.id for e in ctx.engagements.engagements
    }
    for edge in ctx.graph.edges:
        for endpoint in (edge.src, edge.dst):
            if endpoint not in valid:
                yield (
                    f"edge {edge.src} -{edge.kind}-> {edge.dst} has unknown "
                    f"endpoint {endpoint}",
                    "ledger/graph.json",
                )


def graph_04(ctx: Context):
    by_kind: dict[str, int] = {}
    for edge in ctx.graph.edges:
        by_kind[edge.kind] = by_kind.get(edge.kind, 0) + 1

    people = ctx.foundation.people
    externals = ctx.foundation.external_people
    expected = {
        "reports_to": len(people) - 1,
        "works_at": len(people)
        + sum(len(xp.affiliations) or 1 for xp in externals),
        "client_of": len({e.client for e in ctx.engagements.engagements}),
        "participant": sum(
            len(e.internal_participants) + len(e.external_participants)
            for e in ctx.engagements.engagements
        ),
    }
    for kind, want in expected.items():
        got = by_kind.get(kind, 0)
        if got != want:
            yield (
                f"{kind} edge count {got} != {want} derived from ledgers",
                "ledger/graph.json",
            )


# --- LOC (hard-case location policies) --------------------------------------


def _hard_case_hosts(ctx: Context):
    """(entry, key_fact, fact) for every non-body key fact whose fact id
    resolves in the ledger; dangling ids are LOC-03's finding."""
    facts = ctx.engagements.fact_index()
    for e in ctx.manifest:
        for kf in e.key_facts:
            if kf.location != "body" and kf.fact_id in facts:
                yield e, kf, facts[kf.fact_id]


def loc_01(ctx: Context):
    for entry, kf, fact in _hard_case_hosts(ctx):
        if kf.location != "signature_page":
            continue
        if entry.format != "pdf":
            yield (
                f"signature_page fact {fact.id} planted in non-pdf doc",
                entry.path,
            )
            continue
        if not (ctx.paths.share_dir / entry.path).exists():
            continue  # FILE-01/MAN-01 report the absence
        pages = ctx.doc_pages(entry)
        hits = [i for i, text in enumerate(pages) if fact.rendered in text]
        last = len(pages) - 1
        if last not in hits:
            yield (
                f"fact {fact.id} surface {fact.rendered!r} missing from the "
                f"signature page",
                entry.path,
            )
        early = [i + 1 for i in hits if i != last]
        if early:
            yield (
                f"signature-page-only fact {fact.id} leaked into "
                f"page(s) {early}",
                entry.path,
            )


def loc_02(ctx: Context):
    from pathlib import PurePosixPath

    for entry, kf, fact in _hard_case_hosts(ctx):
        if kf.location != "filename":
            continue
        name = PurePosixPath(entry.path).name
        if fact.rendered not in name:
            yield (
                f"filename-only fact {fact.id} surface {fact.rendered!r} "
                f"missing from filename {name!r}",
                entry.path,
            )
        if not (ctx.paths.share_dir / entry.path).exists():
            continue
        if fact.rendered in ctx.doc_text(entry):
            yield (
                f"filename-only fact {fact.id} appears in extractable text",
                entry.path,
            )


def loc_03(ctx: Context):
    facts = ctx.engagements.fact_index()
    hosted: dict[str, int] = {}
    for e in ctx.manifest:
        for kf in e.key_facts:
            fact = facts.get(kf.fact_id)
            if fact is None:
                yield (f"key_fact {kf.fact_id} not in engagement ledger", e.path)
                continue
            if kf.location != fact.location_policy:
                yield (
                    f"key_fact {kf.fact_id} location {kf.location} != ledger "
                    f"policy {fact.location_policy}",
                    e.path,
                )
            if kf.location != "body":
                hosted[kf.fact_id] = hosted.get(kf.fact_id, 0) + 1
    for fid, fact in facts.items():
        if fact.location_policy != "body" and hosted.get(fid, 0) != 1:
            yield (
                f"hard-case fact {fid} ({fact.location_policy}) hosted by "
                f"{hosted.get(fid, 0)} docs, expected exactly 1",
                "docplan/manifest.jsonl",
            )


# --- ACL (read-access overlay) ----------------------------------------------


def _acl_missing(ctx: Context) -> Finding:
    """Any posture other than open requires acl.json; deleting the ledger
    must fail the org, not resurrect the pre-ACL grandfather skip."""
    return (
        f"ledger/acl.json missing but charter posture is "
        f"{ctx.charter.acl_posture!r}",
        "ledger/acl.json",
    )


def acl_01(ctx: Context):
    if ctx.acl is None:
        yield _acl_missing(ctx)
        return
    roster = {p.id for p in ctx.foundation.people}
    for grant in ctx.acl.grants:
        if grant.person not in roster:
            yield (
                f"ACL principal {grant.person} not on the roster",
                "ledger/acl.json",
            )
    for pid in sorted(roster - {g.person for g in ctx.acl.grants}):
        yield (f"roster member {pid} has no ACL grant", "ledger/acl.json")


def acl_02(ctx: Context):
    if ctx.acl is None:
        yield _acl_missing(ctx)
        return
    planned = {e.path for e in ctx.manifest}
    readable: set[str] = set()
    for grant in ctx.acl.grants:
        for doc in grant.docs:
            if doc not in planned:
                yield (
                    f"grant for {grant.person} references unknown doc {doc!r}",
                    "ledger/acl.json",
                )
            readable.add(doc)
    for doc in sorted(planned - readable):
        yield (f"document readable by no one: {doc}", "ledger/acl.json")


def acl_03(ctx: Context):
    if ctx.acl is None:
        yield _acl_missing(ctx)
        return
    from ..acl import derive_acl, render_permissions

    expected = derive_acl(
        ctx.charter, ctx.foundation, ctx.engagements, ctx.manifest
    )
    if expected != ctx.acl:
        yield (
            f"acl.json does not match recomputation from posture "
            f"{ctx.charter.acl_posture!r}",
            "ledger/acl.json",
        )
    if not ctx.paths.permissions_md.exists():
        yield ("PERMISSIONS.md missing from the share", "PERMISSIONS.md")
        return
    # Render from the recomputed ledger (trusted, roster-derived), never
    # from ctx.acl: a tampered grant naming an unknown principal would
    # KeyError inside render_permissions and crash the whole run.
    rendered = render_permissions(ctx.charter, ctx.foundation, expected)
    if ctx.paths.permissions_md.read_text("utf-8") != rendered:
        yield (
            "PERMISSIONS.md does not match the posture's expected rendering",
            "PERMISSIONS.md",
        )


# --- AFF (affiliation-aware documents) --------------------------------------


def _needs_affiliation_docs(ctx: Context) -> str | None:
    if not ctx.charter.graph_targets.affiliations_in_docs:
        return "affiliations_in_docs knob is off for this recipe"
    return None


def aff_01(ctx: Context):
    """Clients and external participants recompute affiliation-aware from
    charter plus foundation. The fabric plan is RNG-free by design, which
    is what makes recomputation tamper evidence."""
    from ..fabric.engagements import affiliation_plan

    range_start = ctx.charter.doc_culture.date_range[0]
    windows = [(e.id, e.start, e.end) for e in ctx.engagements.engagements]
    try:
        plan = affiliation_plan(ctx.foundation, windows, range_start)
    except SystemExit as err:
        yield (
            f"affiliation plan does not recompute from the charter: {err}",
            "foundation.json",
        )
        return
    org_names = {o.id: o.name for o in ctx.foundation.external_orgs}
    for eng in ctx.engagements.engagements:
        client_id, xp_ids = plan[eng.id]
        if eng.client != client_id:
            yield (
                f"client does not recompute affiliation-aware: ledger has "
                f"{eng.client}, recomputation says {client_id}",
                eng.id,
            )
            continue
        if eng.external_participants != xp_ids:
            yield (
                f"external participants do not recompute: ledger has "
                f"{eng.external_participants}, recomputation says {xp_ids}",
                eng.id,
            )
        fact = next(
            (f for f in eng.facts if f.id == f"f:{eng.id}.client"), None
        )
        if fact is None:
            yield (f"engagement has no f:{eng.id}.client fact", eng.id)
        elif fact.rendered != org_names[client_id]:
            yield (
                f"client fact renders {fact.rendered!r}, not the recomputed "
                f"client {org_names[client_id]!r}",
                fact.id,
            )


def aff_02(ctx: Context):
    """Every multi-affiliation person appears in at least one engagement
    per affiliation side (stripped history or stripped participation both
    fire)."""
    from ..fabric.engagements import (
        affiliation_covering,
        padded_window,
        xp_affiliations,
    )

    want = ctx.charter.graph_targets.multi_affiliations
    multi = [
        xp for xp in ctx.foundation.external_people if len(xp.affiliations) >= 2
    ]
    if len(multi) < want:
        yield (
            f"foundation carries {len(multi)} multi-affiliation people, "
            f"charter wants {want}: affiliation history stripped?",
            "foundation.json",
        )
    range_start = ctx.charter.doc_culture.date_range[0]
    for xp in multi:
        for aff in xp_affiliations(xp):
            on_side = any(
                eng.client == aff.org
                and xp.id in eng.external_participants
                and affiliation_covering(
                    xp,
                    aff.org,
                    *padded_window(eng.start, eng.end, range_start),
                )
                for eng in ctx.engagements.engagements
            )
            if not on_side:
                yield (
                    f"{xp.id} never participates in an engagement under "
                    f"{aff.org} inside that affiliation",
                    "ledger/engagements.json",
                )


# --- EML ------------------------------------------------------------------


def eml_01(ctx: Context):
    """Transport headers recompute exactly from the ledgers, via the same
    helper the renderer used. Runs only when the charter asks for mail, so
    a knob that is on with no .eml documents is tamper evidence."""
    from ..render import people_index
    from ..render.eml import expected_headers, thread_members

    # Derived noise .eml files (M12) mirror the source they copy or draft;
    # their headers are not independently recomputable from the ledger (an
    # exact duplicate carries the source's headers verbatim), so they are
    # excluded here, exactly as SCAN-01/LEG-01 exclude derived docs.
    entries = [
        e
        for e in ctx.manifest
        if e.format == "eml" and e.authoring != "derived"
    ]
    if not entries:
        yield (
            "format_mix.eml > 0 but the manifest plans no eml documents",
            "docplan/manifest.jsonl",
        )
        return
    people = people_index(ctx.foundation)
    for e in entries:
        path = ctx.paths.share_dir / e.path
        if not path.exists():
            continue  # FILE-01/MAN-01 report the absence
        msg = _eml_message(path)
        thread = thread_members(e, ctx.manifest)
        expected = expected_headers(
            e, people, ctx.charter.slug, ctx.charter.domain, thread
        )
        for name, want in expected.items():
            got = re.sub(r"\s+", " ", str(msg[name] or "")).strip()
            if got != re.sub(r"\s+", " ", want).strip():
                yield (
                    f"header {name} does not recompute from the ledger: "
                    f"{got!r} != {want!r}",
                    e.path,
                )
        # A thread opener must carry no threading headers: a planted
        # In-Reply-To / References on an opener is tamper evidence the
        # per-header loop above cannot see (it only checks expected headers).
        if "In-Reply-To" not in expected:
            for spurious in ("In-Reply-To", "References"):
                if msg[spurious] is not None:
                    yield (
                        f"header {spurious} present on a message the ledger "
                        f"gives no thread predecessor",
                        e.path,
                    )


def eml_02(ctx: Context):
    """Mail-block messages end their own words with a signature block that
    recomputes from the ledger (name, title AS OF the send date, phone), so a
    promotion changes it mid-corpus and the model cannot author it. Runs only
    when the recipe declares doc_culture.mail; a knob on with no mail is tamper
    evidence, exactly like EML-01."""
    from ..render.eml import mail_signature

    entries = [
        e
        for e in ctx.manifest
        if e.format == "eml"
        and e.authoring != "derived"
        and "send_minute" in e.render_params
    ]
    if not entries:
        yield (
            "doc_culture.mail is declared but the manifest plans no "
            "thread mail",
            "docplan/manifest.jsonl",
        )
        return
    for e in entries:
        path = ctx.paths.share_dir / e.path
        if not path.exists():
            continue  # FILE-01/MAN-01 report the absence
        body = _eml_message(path).get_body(preferencelist=("plain",))
        text = body.get_content() if body is not None else ""
        want = mail_signature(ctx.foundation.person(e.authors[0]), e.date)
        if want not in text:
            yield (
                "mail signature block does not recompute from the ledger "
                f"(expected the author's name / title-as-of-{e.date} / phone)",
                e.path,
            )


def eml_03(ctx: Context):
    """A transmittal email's MIME attachment is byte-identical to the share
    document it carries (render_params.attach_path). Runs when mail is
    declared; a recipe asking for attachments with none planted, or a
    transmittal whose attachment does not match, is a failure, not a skip."""
    from ..render.eml import eml_attachment_bytes

    transmittals = [
        e
        for e in ctx.manifest
        if e.render_params.get("attach_path") and e.authoring != "derived"
    ]
    want = ctx.charter.doc_culture.mail.attachments
    if want > 0 and not transmittals:
        yield (
            f"mail.attachments is {want} but the manifest plans no "
            "transmittal mail",
            "docplan/manifest.jsonl",
        )
        return
    for e in transmittals:
        path = ctx.paths.share_dir / e.path
        src = ctx.paths.share_dir / str(e.render_params["attach_path"])
        if not path.exists() or not src.exists():
            continue  # FILE-01/MAN-01 report the absence
        got = eml_attachment_bytes(path)
        if got is None:
            yield ("transmittal email carries no MIME attachment", e.path)
        elif got != src.read_bytes():
            yield (
                "transmittal attachment is not byte-identical to "
                f"{e.render_params['attach_path']!r}",
                e.path,
            )


def dl_01(ctx: Context):
    """The distribution-list ledger recomputes from charter + foundation;
    every DL-addressed message names a real list; and visibility expands the
    list to its members (every current member can read the message). Runs when
    the recipe declares lists; a knob on with the ledger missing is tamper
    evidence, not a skip."""
    from ..acl import derive_distribution_lists

    expected = derive_distribution_lists(ctx.charter, ctx.foundation)
    if ctx.dls is None:
        yield (
            "distribution_lists declared but ledger/distribution_lists.json "
            "is missing",
            "ledger/distribution_lists.json",
        )
        return
    if ctx.dls.model_dump() != expected.model_dump():
        yield (
            "distribution-list ledger does not recompute from charter + "
            "foundation",
            "ledger/distribution_lists.json",
        )
        return
    by_addr = {dl.address: dl for dl in ctx.dls.lists}
    current = {
        p.id for p in ctx.foundation.people if p.employment.end is None
    }
    reader_of = {}
    if ctx.acl is not None:
        for grant in ctx.acl.grants:
            for doc in grant.docs:
                reader_of.setdefault(doc, set()).add(grant.person)
    for e in ctx.manifest:
        addr = e.render_params.get("dl")
        if not addr:
            continue
        dl = by_addr.get(str(addr))
        if dl is None:
            yield (
                f"message addresses {addr!r}, not a known distribution list",
                e.path,
            )
            continue
        if ctx.acl is None:
            continue  # visibility expansion is ACL-derived; ACL-* report absence
        missing = (set(dl.members) & current) - reader_of.get(e.path, set())
        if missing:
            yield (
                f"distribution list {dl.address} members {sorted(missing)} "
                f"cannot read the message addressed to them",
                e.path,
            )


# --- SCAN -----------------------------------------------------------------


def scan_01(ctx: Context):
    """Scan flags recompute from the charter, scanned docs are raster
    pdfs, and extractable text matches the planned OCR-layer presence."""
    import pikepdf
    from pypdf import PdfReader

    from ..docplan.planner import scan_selection

    policy = {
        fid: f.location_policy
        for fid, f in ctx.engagements.fact_index().items()
    }
    expected = scan_selection(
        ctx.charter.seed,
        ctx.charter.doc_culture,
        [
            (
                e.date,
                e.path,
                any(
                    policy.get(ref) == "signature_page" for ref in e.facts_refs
                ),
            )
            for e in ctx.manifest
            if e.format == "pdf" and e.authoring != "derived"
        ],
    )
    for e in ctx.manifest:
        if e.format != "pdf" or e.authoring == "derived":
            # Derived noise docs (M12) are never scanned: they are excluded
            # from the recompute, so a clean copy or draft of a pdf does not
            # perturb the scan-set ordering.
            continue
        want_layer = expected.get(e.path)  # None=plain, False=image, True=ocr
        got_scan = e.render_params.get("scan") == 1
        got_layer = e.render_params.get("ocr_layer") == 1
        if got_scan != (want_layer is not None) or got_layer != bool(want_layer):
            yield (
                f"scan flags do not recompute from the charter: manifest has "
                f"scan={got_scan} ocr_layer={got_layer}, expected "
                f"scan={want_layer is not None} ocr_layer={bool(want_layer)}",
                e.path,
            )
            continue
        if want_layer is None:
            continue
        path = ctx.paths.share_dir / e.path
        if not path.exists():
            continue  # FILE-01/MAN-01 report the absence
        try:
            with pikepdf.open(path) as pdf:
                for page in pdf.pages:
                    xobjects = page.get("/Resources", {}).get("/XObject", {})
                    if not any(
                        x.get("/Subtype") == pikepdf.Name.Image
                        for x in xobjects.values()
                    ):
                        yield ("scanned doc has a page with no raster image",
                               e.path)
                        break
        except pikepdf.PdfError as err:
            yield (f"scanned doc does not open: {err}", e.path)
            continue
        text = "\n".join(
            page.extract_text() or "" for page in PdfReader(str(path)).pages
        )
        if want_layer and not text.strip():
            yield ("planned OCR layer exposes no extractable text", e.path)
        if not want_layer and text.strip():
            yield ("image-only scan exposes extractable text", e.path)


def scan_02(ctx: Context):
    """The true-text archive exists exactly for scanned docs and ties to
    the rendered page count. Image-only text obligations run against this
    archive via doc_text/doc_pages, so FACT/MENT/LOC own those findings."""
    from pypdf import PdfReader

    from ..render.scan import scan_pages_path

    scanned = {
        scan_pages_path(ctx.paths, e.doc_id).name: e
        for e in ctx.manifest
        if e.render_params.get("scan") == 1
    }
    on_disk = (
        {p.name for p in ctx.paths.scans_dir.glob("*.pages.json")}
        if ctx.paths.scans_dir.exists()
        else set()
    )
    for name in sorted(set(scanned) - on_disk):
        yield ("scanned doc has no archived page text", f"scans/{name}")
    for name in sorted(on_disk - set(scanned)):
        yield ("archived page text for a doc the plan never scanned",
               f"scans/{name}")
    for name, entry in sorted(scanned.items()):
        if name not in on_disk:
            continue
        pages = ctx.scan_archive(entry)
        if pages is None:
            yield ("archived page text does not parse", f"scans/{name}")
            continue
        rendered = ctx.paths.share_dir / entry.path
        if not rendered.exists():
            continue  # FILE-01/MAN-01 report the absence
        try:
            count = len(PdfReader(str(rendered)).pages)
        except Exception:  # noqa: BLE001 - unreadable = finding
            yield ("rendered pdf does not open", entry.path)
            continue
        if len(pages) != count:
            yield (
                f"archive holds {len(pages)} page(s), rendered pdf has "
                f"{count}",
                entry.path,
            )


# --- LEG ------------------------------------------------------------------


def leg_01(ctx: Context):
    """Legacy assignment recomputes from the charter, and every legacy
    binary is an OLE container holding the stream its format promises."""
    import olefile

    from ..docplan.planner import legacy_selection
    from ..render.legacy import LEGACY_STREAMS
    from ..schemas import BASE_FORMAT

    def modern_path(e) -> str:
        if e.format in BASE_FORMAT:
            return e.path[: -len(e.format)] + BASE_FORMAT[e.format]
        return e.path

    expected = legacy_selection(
        ctx.charter.doc_culture,
        [
            (e.date, modern_path(e))
            for e in ctx.manifest
            if (e.format in BASE_FORMAT or e.format in set(BASE_FORMAT.values()))
            and e.authoring != "derived"
        ],
    )
    for e in ctx.manifest:
        if e.authoring == "derived":
            # Derived noise docs (M12) are never legacy binaries: excluded
            # from the recompute so a modern-format copy does not shift the
            # oldest-N legacy ordering.
            continue
        is_legacy = e.format in BASE_FORMAT
        if not is_legacy and e.format not in set(BASE_FORMAT.values()):
            continue
        should_be = modern_path(e) in expected
        if is_legacy != should_be:
            yield (
                f"legacy assignment does not recompute from the charter: "
                f"manifest format is {e.format!r}, expected "
                f"{'legacy' if should_be else 'modern'}",
                e.path,
            )
            continue
        if not is_legacy:
            continue
        path = ctx.paths.share_dir / e.path
        if not path.exists():
            continue  # FILE-01/MAN-01 report the absence
        try:
            is_ole = olefile.isOleFile(str(path))
        except Exception:  # noqa: BLE001 - unreadable = finding
            is_ole = False
        if not is_ole:
            yield ("legacy file is not an OLE container", e.path)
            continue
        stream = LEGACY_STREAMS[e.format]
        try:
            with olefile.OleFileIO(str(path)) as ole:
                has_stream = ole.exists(stream)
        except Exception:  # noqa: BLE001 - unreadable = finding
            yield ("legacy file does not open as an OLE container", e.path)
            continue
        if not has_stream:
            yield (
                f"OLE container has no {stream!r} stream; not a real "
                f".{e.format}",
                e.path,
            )


# --- PROV -----------------------------------------------------------------


def prov_01(ctx: Context):
    from ..render.provenance import (
        eml_has_marker,
        legacy_has_marker,
        opc_has_marker,
        pdf_has_marker,
    )

    checkers = {"docx": opc_has_marker, "pdf": pdf_has_marker,
                "xlsx": opc_has_marker, "pptx": opc_has_marker,
                "eml": eml_has_marker, "doc": legacy_has_marker,
                "xls": legacy_has_marker, "ppt": legacy_has_marker}
    for e in ctx.manifest:
        path = ctx.paths.share_dir / e.path
        if not path.exists():
            continue
        checker = checkers.get(e.format)
        if checker is None:
            # A format without a marker checker can never silently pass.
            yield (f"no provenance checker known for format {e.format!r}", e.path)
        elif not checker(path):
            yield ("synthetic-provenance marker missing", e.path)


RULES = [
    Rule("ORG-01", "ERROR", "exactly one CEO-equivalent", org_01),
    Rule("ORG-02", "ERROR", "reports_to is a single acyclic tree", org_02),
    Rule("NAME-01", "ERROR", "no name collides with a known real firm",
         name_01),
    Rule("DATE-01", "ERROR", "doc dates inside charter range", date_01),
    Rule("DATE-02", "ERROR", "authors employed at doc date", date_02),
    Rule("CAL-01", "ERROR", "attendance genres land on business days", cal_01,
         available=_needs_calendar),
    Rule("NOISE-01", "ERROR", "derived noise docs name a real source; "
         "version chains cohere and diverge",
         noise_01, available=_needs_noise),
    Rule("FIN-01", "ERROR", "finance ledger ties out", fin_01),
    Rule("FIN-02", "ERROR", "workbooks tie to the finance ledger", fin_02),
    Rule("FACT-01", "ERROR", "planted facts appear verbatim in doc text", fact_01),
    Rule("MENT-01", "ERROR", "planned mention surfaces appear in doc text",
         ment_01, available=_needs_mentions),
    Rule("MENT-02", "ERROR", "mentions resolve to graph entities", ment_02,
         available=_needs_mentions),
    Rule("GRAPH-01", "ERROR", "per-person mention coverage meets the recipe",
         graph_01, available=_needs_mention_knob),
    Rule("GRAPH-02", "ERROR", "no orphan roster member", graph_02,
         available=_needs_mention_knob),
    Rule("GRAPH-03", "ERROR", "graph edges have known endpoints", graph_03),
    Rule("GRAPH-04", "ERROR", "per-type edge counts match the ledgers",
         graph_04),
    Rule("LOC-01", "ERROR", "signature-page facts on the signature page only",
         loc_01),
    Rule("LOC-02", "ERROR", "filename facts in the filename only", loc_02),
    Rule("LOC-03", "ERROR", "manifest key-fact locations mirror the ledger",
         loc_03),
    Rule("ACL-01", "ERROR", "ACL principals mirror the roster", acl_01,
         available=_needs_acl),
    Rule("ACL-02", "ERROR", "every document is readable by someone", acl_02,
         available=_needs_acl),
    Rule("ACL-03", "ERROR", "grants and PERMISSIONS.md match the posture",
         acl_03, available=_needs_acl),
    Rule("AFF-01", "ERROR", "clients and external participants recompute "
         "affiliation-aware", aff_01, available=_needs_affiliation_docs),
    Rule("AFF-02", "ERROR", "multi-affiliation people appear under both "
         "employers", aff_02, available=_needs_affiliation_docs),
    Rule("EML-01", "ERROR", "eml transport headers recompute from the ledger",
         eml_01, available=_needs_eml),
    Rule("EML-02", "ERROR", "mail signature blocks recompute from the ledger",
         eml_02, available=_needs_mail),
    Rule("EML-03", "ERROR", "transmittal attachments match their share doc",
         eml_03, available=_needs_mail),
    Rule("DL-01", "ERROR", "distribution lists recompute and expand for "
         "visibility", dl_01, available=_needs_dls),
    Rule("STY-01", "ERROR", "per-person style specs recompute from the "
         "roster", sty_01, available=_needs_style),
    Rule("SCAN-01", "ERROR", "scan flags recompute; raster and OCR presence "
         "match the plan", scan_01, available=_needs_scan),
    Rule("SCAN-02", "ERROR", "true-text archives exist exactly for scans",
         scan_02, available=_needs_scan),
    Rule("LEG-01", "ERROR", "legacy assignment recomputes; binaries are real "
         "OLE containers", leg_01, available=_needs_legacy),
    Rule("FILE-01", "ERROR", "every manifest doc opens in its native reader",
         file_01),
    Rule("MAN-01", "ERROR", "manifest and share tree match 1:1", man_01),
    Rule("PROV-01", "ERROR", "every rendered file carries the synthetic marker",
         prov_01),
]
